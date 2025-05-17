
import streamlit as st
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO

# Branding colors
SALVATION_RED = "#c8102e"
HEADER_HTML = f"""
<div style='background-color:{SALVATION_RED}; padding: 1rem; border-radius: 0 0 6px 6px; text-align: center;'>
    <h1 style='color: white; font-family: Arial, sans-serif; margin: 0;'>SASB Song Combiner</h1>
    <p style='color: white; font-size: 1rem; margin-top: 0.2rem;'>A tool inspired by the mission of The Salvation Army</p>
</div>
"""

st.set_page_config(page_title="SASB Song Combiner", layout="centered")

# Custom header
st.markdown(HEADER_HTML, unsafe_allow_html=True)
st.markdown("### 🎶 Quickly generate a PowerPoint slideshow of SASB songs", unsafe_allow_html=True)

# Styled form
with st.form("song_form"):
    st.markdown("#### ✏️ Song List", unsafe_allow_html=True)
    song_nums = st.text_input("Enter Song Numbers (comma-separated)", "2, 3, 5")

    st.markdown("#### 🎨 Appearance Options", unsafe_allow_html=True)
    font_color = st.color_picker("Font Color", "#FFFFFF")
    bg_color = st.color_picker("Background Color", "#000000")
    bg_image = st.file_uploader("Upload Background Image (optional)", type=["png", "jpg", "jpeg"])
    logo_image = st.file_uploader("Upload Logo (optional)", type=["png", "jpg", "jpeg"])

    submit = st.form_submit_button("📄 Generate PowerPoint")

def split_text_into_chunks(lines, chunk_size):
    return [lines[i:i + chunk_size] for i in range(0, len(lines), chunk_size)]

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def create_combined_pptx(song_numbers, font_color_hex, bg_color_hex, bg_img_bytes, logo_bytes):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    if not os.path.exists("songs"):
        st.error("⚠️ The 'songs' folder is missing. Please upload .pptx files to a folder named 'songs'.")
        st.stop()

    for num in song_numbers:
        match = next((f for f in os.listdir("songs") if f.startswith(f"{num} ")), None)
        if not match:
            continue

        src = Presentation(os.path.join("songs", match))
        for slide in src.slides:
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        lines.append(para.text.strip())
            if not lines:
                continue

            footer = lines[-1] if len(lines) > 1 else ""
            chunks = split_text_into_chunks(lines[:-1], 8)

            for chunk in chunks:
                s = prs.slides.add_slide(prs.slide_layouts[6])

                # Background
                s.background.fill.solid()
                r, g, b = hex_to_rgb(bg_color_hex)
                s.background.fill.fore_color.rgb = RGBColor(r, g, b)

                if bg_img_bytes:
                    s.shapes.add_picture(bg_img_bytes, 0, 0, width=prs.slide_width, height=prs.slide_height)

                # Main text box
                textbox = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(5.5))
                tf = textbox.text_frame
                tf.clear()
                tf.vertical_anchor = 1
                tf.word_wrap = True

                para = tf.paragraphs[0]
                para.text = "\n".join(chunk)
                para.alignment = 1
                run = para.runs[0]
                run.font.size = Pt(44)
                run.font.name = 'Calibri'
                run.font.bold = True
                run.font.color.rgb = RGBColor(*hex_to_rgb(font_color_hex))

                # Footer
                footer_box = s.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.5))
                footer_tf = footer_box.text_frame
                p = footer_tf.paragraphs[0]
                p.clear()
                run = p.add_run()
                run.text = footer
                run.font.size = Pt(20)
                run.font.name = "Calibri"
                run.font.bold = True
                run.font.color.rgb = RGBColor(*hex_to_rgb(font_color_hex))
                p.alignment = 1

                if logo_bytes:
                    s.shapes.add_picture(logo_bytes, Inches(0.2), Inches(6.4), width=Inches(1.0))

        # Add blank slide
        blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
        blank_slide.background.fill.solid()
        r, g, b = hex_to_rgb(bg_color_hex)
        blank_slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

if submit:
    song_list = [s.strip() for s in song_nums.split(",") if s.strip().isdigit()]
    bg_img_bytes = bg_image if bg_image is not None else None
    logo_bytes = logo_image if logo_image is not None else None
    pptx_file = create_combined_pptx(song_list, font_color, bg_color, bg_img_bytes, logo_bytes)

    st.success("✅ PowerPoint created!")
    st.download_button("⬇️ Download Presentation", pptx_file, file_name="combined_songs.pptx")
